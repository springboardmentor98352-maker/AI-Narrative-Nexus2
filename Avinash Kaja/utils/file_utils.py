# utils/file_utils.py

import pandas as pd
import json
from bs4 import BeautifulSoup
from docx import Document
import pdfplumber
from pptx import Presentation
import striprtf
from lxml import etree as ET


# PDF HANDLER
def read_pdf(file):
    try:
        text = []
        with pdfplumber.open(file) as pdf:
            for page in pdf.pages:
                extracted = page.extract_text() or ""
                text.append(extracted)
        return "\n".join(text)
    except Exception as e:
        return f"PDF read error: {e}"


# DOCX HANDLER
def read_docx(file):
    try:
        doc = Document(file)
        return "\n".join([p.text for p in doc.paragraphs])
    except Exception as e:
        return f"DOCX read error: {e}"


# TXT HANDLER
# TXT HANDLER
def read_txt(file):
    try:
        return file.read().decode("utf-8", errors="ignore")
    except Exception:
        try:
            return file.read().decode("latin-1")
        except Exception:
            return "TXT decoding failed."


# PPTX HANDLER
def read_pptx(file):
    try:
        text = []
        prs = Presentation(file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        return f"PPTX read error: {e}"


# CSV HANDLER
def read_csv(file):
    try:
        df = pd.read_csv(file)
        return "\n".join(df.astype(str).agg(" ".join, axis=1))
    except Exception as e:
        return f"CSV read error: {e}"


# EXCEL HANDLER
def read_excel(file):
    try:
        df = pd.read_excel(file)
        return "\n".join(df.astype(str).agg(" ".join, axis=1))
    except Exception as e:
        return f"Excel read error: {e}"


# JSON HANDLER
def read_json(file):
    try:
        data = json.load(file)
        return json.dumps(data, indent=2)
    except Exception as e:
        return f"JSON read error: {e}"


# HTML HANDLER
def read_html(file):
    try:
        soup = BeautifulSoup(file.read(), "html.parser")
        return soup.get_text(separator="\n")
    except Exception as e:
        return f"HTML read error: {e}"


# RTF HANDLER
def read_rtf(file):
    try:
        rtf_content = file.read().decode("latin-1")
        text = striprtf.rtf_to_text(rtf_content)
        return text
    except Exception as e:
        return f"RTF read error: {e}"


# EPUB HANDLER
def read_epub(file):
    try:
        raw = file.read().decode("utf-8", errors="ignore")
        soup = BeautifulSoup(raw, "html.parser")
        return soup.get_text(separator="\n")
    except Exception as e:
        return f"EPUB read error: {e}"


# XML HANDLER
def read_xml(file):
    try:
        tree = ET.parse(file)
        root = tree.getroot()
        return ET.tostring(root, encoding="unicode")
    except Exception as e:
        return f"XML read error: {e}"


# MASTER HANDLER
def read_file(file):
    name = file.name.lower()

    if name.endswith(".pdf"):
        return read_pdf(file)
    elif name.endswith(".docx"):
        return read_docx(file)
    elif name.endswith(".txt"):
        return read_txt(file)
    elif name.endswith(".pptx"):
        return read_pptx(file)
    elif name.endswith(".csv"):
        return read_csv(file)
    elif name.endswith(".xlsx"):
        return read_excel(file)
    elif name.endswith(".json"):
        return read_json(file)
    elif name.endswith(".html") or name.endswith(".htm"):
        return read_html(file)
    elif name.endswith(".rtf"):
        return read_rtf(file)
    elif name.endswith(".epub"):
        return read_epub(file)
    elif name.endswith(".xml"):
        return read_xml(file)
    else:
        return "Unsupported file format."
